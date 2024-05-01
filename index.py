from pptx import Presentation
from PIL import Image
import io
import os

def resize_image(image_stream):
  first_output_width = 600
  second_out_put_width = 300
  image = Image.open(image_stream)
  original_format = image.format

  if image.format == 'WMF':
    image = image.convert('RGB')  # WMFはサポート外なのでRGBに変換
    original_format = 'PNG'       # PNG形式で保存
  print(image.width)
  if image.width > first_output_width:
    aspect_ratio = image.height / image.width
    new_height = int(first_output_width * aspect_ratio)
    resized_image = image.resize((first_output_width, new_height), Image.Resampling.LANCZOS)
    output_stream = io.BytesIO()
    resized_image.save(output_stream, format=original_format)
    output_stream.seek(0)
    return output_stream
  elif image.width > second_out_put_width:
    aspect_ratio = image.height / image.width
    new_height = int(second_out_put_width * aspect_ratio)
    resized_image = image.resize((second_out_put_width, new_height), Image.Resampling.LANCZOS)
    output_stream = io.BytesIO()
    resized_image.save(output_stream, format=original_format)
    output_stream.seek(0)
    return output_stream
  else:
    print('not resized')
    output_stream = io.BytesIO()
    image.save(output_stream, format=original_format)
    output_stream.seek(0)
    return output_stream


def replace_image_in_pptx(pptx_path, output_path):
  prs = Presentation(pptx_path)
  
  for slide in prs.slides:
    for shape in list(slide.shapes):  # シェイプリストをイテレート
      if shape.shape_type == 13:  # 画像の場合のみ処理
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        image_stream = io.BytesIO(shape.image.blob)
        new_image_stream = resize_image(image_stream)
        sp = shape._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(new_image_stream, left, top, width, height)  # 新しいJPEG画像を挿入

  prs.save(output_path)

if __name__ == '__main__':
  import glob
  files = glob.iglob('./data/*.pptx')
  
  for file_path in files:
    file_name = os.path.basename(file_path)
    output_folder = './out'
    output_pptx_path = f'{output_folder}/{file_name}'
    os.makedirs(output_folder, exist_ok=True)
    replace_image_in_pptx(file_path, output_pptx_path)
